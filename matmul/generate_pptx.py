#!/usr/bin/env python3
"""Generate a PowerPoint presentation about GPU Matrix Multiplication."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Helpers ──────────────────────────────────────────────────────────────────

def set_slide_bg(slide, r, g, b):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=RGBColor(0x33, 0x33, 0x33),
                alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide_content(slide, left, top, width, height, items,
                              font_size=18, color=RGBColor(0x33, 0x33, 0x33),
                              font_name="Calibri", spacing=Pt(12)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
    return txBox

def add_code_box(slide, left, top, width, height, code_text, font_size=13):
    """Add a dark code box with monospace text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)  # dark slate
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)

    for i, line in enumerate(code_text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        p.font.name = "Courier New"
        p.space_after = Pt(2)
    return shape

def add_colored_rect(slide, left, top, width, height, fill_rgb, text="",
                     font_size=12, font_color=RGBColor(0xFF, 0xFF, 0xFF)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape

def add_section_header(slide, text):
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.7),
                text, font_size=32, bold=True, color=RGBColor(0x1E, 0x3A, 0x5F))
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.05), Inches(1.5), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x3B, 0x82, 0xF6)
    line.line.fill.background()

def draw_matrix_grid(slide, left, top, cell_size, rows, cols, fill_rgb,
                      label="", highlight_cells=None, highlight_rgb=None):
    """Draw a matrix grid of colored squares."""
    if label:
        add_textbox(slide, left, top - Inches(0.35), Inches(cols * cell_size),
                    Inches(0.3), label, font_size=11, bold=True,
                    color=fill_rgb, alignment=PP_ALIGN.CENTER)

    for r in range(rows):
        for c in range(cols):
            x = left + Emu(int(c * Inches(cell_size)))
            y = top + Emu(int(r * Inches(cell_size)))
            w = Inches(cell_size * 0.92)
            h = Inches(cell_size * 0.92)

            is_highlighted = highlight_cells and (r, c) in highlight_cells
            rgb = highlight_rgb if is_highlighted else fill_rgb

            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb
            shape.line.fill.background()
            if is_highlighted:
                shape.shadow.inherit = False


# ── Presentation ─────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = (0x0F, 0x17, 0x2A)
WHITE_BG = (0xF8, 0xFA, 0xFC)
LIGHT_BG = (0xF1, 0xF5, 0xF9)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1 – Title
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, *DARK_BG)

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
            "Matrix Multiplication", font_size=52, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF))

add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(1.0),
            "Down the Rabbit Hole", font_size=40, bold=False,
            color=RGBColor(0x93, 0xC5, 0xFD))

# Accent line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(1), Inches(4.2), Inches(3), Inches(0.06))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0x3B, 0x82, 0xF6)
line.line.fill.background()

add_textbox(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.6),
            "A journey through GPU GEMM optimization — from naive to blazing fast",
            font_size=20, color=RGBColor(0x94, 0xA3, 0xB8))

add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
            "SGEMM  |  CUDA  |  Memory Coalescing  |  Shared Memory  |  Tensor Cores",
            font_size=14, color=RGBColor(0x64, 0x74, 0x8B))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2 – What is GEMM?
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, *WHITE_BG)
add_section_header(slide, "What is GEMM?")

add_textbox(slide, Inches(0.8), Inches(1.4), Inches(7), Inches(0.6),
            "General Matrix Multiply:  C = \u03b1(A \u00d7 B) + \u03b2C",
            font_size=28, bold=True, color=RGBColor(0x1E, 0x40, 0xAF),
            font_name="Calibri")

bullets = [
    "\u2022  The workhorse of deep learning — every linear layer, attention head, and convolution is a GEMM",
    "\u2022  A is M\u00d7K, B is K\u00d7N, C is M\u00d7N — billions of multiply-add operations",
    "\u2022  GPUs have thousands of cores — perfect for this embarrassingly parallel problem",
    "\u2022  But getting peak performance requires understanding the hardware deeply",
]
add_bullet_slide_content(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(3.5),
                          bullets, font_size=20, spacing=Pt(16))

# Formula box
add_code_box(slide, Inches(0.8), Inches(5.0), Inches(6), Inches(1.2),
             "for (int k = 0; k < K; ++k)\n"
             "    C[i][j] += A[i][k] * B[k][j];",
             font_size=18)

add_textbox(slide, Inches(7.2), Inches(5.2), Inches(5), Inches(0.8),
            "Each thread computes one element of C\n"
            "by walking along the shared K dimension",
            font_size=16, color=RGBColor(0x64, 0x74, 0x8B))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3 – GPU Execution Model
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, *WHITE_BG)
add_section_header(slide, "GPU Execution Model: Threads, Warps & Memory")

# Left column — Thread hierarchy
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
            "Thread Hierarchy", font_size=20, bold=True,
            color=RGBColor(0x1E, 0x3A, 0x5F))

hierarchy_items = [
    "\u2022  Grid \u2192 Blocks \u2192 Threads",
    "\u2022  Threads are grouped into Warps (32 threads)",
    "\u2022  All threads in a warp execute in lockstep (SIMT)",
    "\u2022  threadIdx.x changes fastest across a warp",
    "\u2022  Adjacent threads = consecutive threadIdx.x values",
]
add_bullet_slide_content(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3.0),
                          hierarchy_items, font_size=18, spacing=Pt(14))

# Right column — Memory hierarchy
add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
            "Memory Hierarchy", font_size=20, bold=True,
            color=RGBColor(0x1E, 0x3A, 0x5F))

mem_items = [
    "\u2022  Registers — per thread, fastest (~0 cycles)",
    "\u2022  Shared Memory — per block, ~5 cycles",
    "\u2022  Global Memory (DRAM) — ~400-600 cycles",
    "\u2022  DRAM serves data in 32-byte transactions",
    "\u2022  Wasted bytes in a transaction = wasted bandwidth",
]
add_bullet_slide_content(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(3.0),
                          mem_items, font_size=18, spacing=Pt(14))

# Key insight box
box = add_colored_rect(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.2),
                        RGBColor(0xEF, 0xF6, 0xFF))
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
p = tf.paragraphs[0]
p.text = ("\U0001f4a1  Key Insight: When adjacent threads (in the same warp) access adjacent "
          "memory addresses, the GPU can combine these into a single DRAM transaction. "
          "This is called memory coalescing — it is the #1 performance factor for "
          "memory-bound kernels.")
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)
p.font.name = "Calibri"

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4 – The Naive Kernel (Anti-Pattern)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, *WHITE_BG)
add_section_header(slide, "Stage 0: The Naive Kernel (Anti-Pattern)")

add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
            "The most intuitive mapping — but it destroys memory bandwidth",
            font_size=18, color=RGBColor(0xDC, 0x26, 0x26))

kernel_code = (
    "__global__ void sgemm_uncoalesced(\n"
    "    int M, int N, int K, float alpha,\n"
    "    const float *A, const float *B,\n"
    "    float beta, float *C)\n"
    "{\n"
    "  // threadIdx.x -> row (i), threadIdx.y -> col (j)\n"
    "  const uint i = blockIdx.x * blockDim.x + threadIdx.x;\n"
    "  const uint j = blockIdx.y * blockDim.y + threadIdx.y;\n"
    "\n"
    "  if (i < M && j < N) {\n"
    "    float tmp = 0.0;\n"
    "    for (int k = 0; k < K; ++k)\n"
    "      tmp += A[i * K + k] * B[k * N + j];\n"
    "    C[i * N + j] = alpha * tmp + beta * C[i * N + j];\n"
    "  }\n"
    "}"
)
add_code_box(slide, Inches(0.8), Inches(2.0), Inches(7), Inches(4.2),
             kernel_code, font_size=14)

# Annotations on the right
add_textbox(slide, Inches(8.2), Inches(2.0), Inches(4.5), Inches(0.4),
            "What goes wrong?", font_size=20, bold=True,
            color=RGBColor(0xDC, 0x26, 0x26))

right_items = [
    "\u2022  threadIdx.x maps to row i",
    "\u2022  Adjacent threads in a warp\n   have different i values",
    "\u2022  A[i*K+k]: adjacent threads\n   access different rows \u2192 strided!",
    "\u2022  B[k*N+j]: j is constant within\n   the warp \u2192 broadcast (1 value)",
    "\u2022  C[i*N+j]: writes are strided too",
]
add_bullet_slide_content(slide, Inches(8.2), Inches(2.6), Inches(4.5), Inches(3.5),
                          right_items, font_size=16, spacing=Pt(14),
                          color=RGBColor(0x47, 0x47, 0x47))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5 – Row-Major Memory Layout
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, *WHITE_BG)
add_section_header(slide, "Row-Major Memory Layout")

add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.6),
            "A 2D matrix is stored as a flat 1D array in DRAM — row after row",
            font_size=20, color=RGBColor(0x47, 0x47, 0x47))

# Show 4x4 matrix conceptually
add_textbox(slide, Inches(0.8), Inches(2.2), Inches(3), Inches(0.4),
            "Logical 2D View (4\u00d74):", font_size=16, bold=True,
            color=RGBColor(0x1E, 0x3A, 0x5F))

colors_2d = [
    RGBColor(0xEF, 0x44, 0x44),  # row 0 - red
    RGBColor(0x22, 0xC5, 0x5E),  # row 1 - green
    RGBColor(0x3B, 0x82, 0xF6),  # row 2 - blue
    RGBColor(0xF5, 0x9E, 0x0B),  # row 3 - amber
]

cell = 0.5
for r in range(4):
    for c in range(4):
        x = Inches(1.0 + c * 0.6)
        y = Inches(2.7 + r * 0.6)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x, y, Inches(0.52), Inches(0.52))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors_2d[r]
        shape.line.fill.background()
        tf = shape.text_frame
        tf.margin_left = Pt(0)
        tf.margin_right = Pt(0)
        p = tf.paragraphs[0]
        p.text = f"[{r},{c}]"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

# Arrow
add_textbox(slide, Inches(3.8), Inches(3.5), Inches(1), Inches(0.5),
            "\u27a1", font_size=36, color=RGBColor(0x64, 0x74, 0x8B),
            alignment=PP_ALIGN.CENTER)

# Show 1D DRAM layout
add_textbox(slide, Inches(5.0), Inches(2.2), Inches(7), Inches(0.4),
            "Physical 1D Layout in DRAM:", font_size=16, bold=True,
            color=RGBColor(0x1E, 0x3A, 0x5F))

for i in range(16):
    row = i // 4
    col = i % 4
    x = Inches(5.0 + i * 0.48)
    y = Inches(2.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x, y, Inches(0.44), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors_2d[row]
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    p = tf.paragraphs[0]
    p.text = f"{i}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# Address labels
add_textbox(slide, Inches(5.0), Inches(3.45), Inches(7.7), Inches(0.3),
            "addr: 0  1  2  3  4  5  6  7  8  9  10 11 12 13 14 15",
            font_size=10, color=RGBColor(0x94, 0xA3, 0xB8), font_name="Courier New")

# Explanation
add_textbox(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.4),
            "Indexing formula:   A[row][col]  \u2192  A[row * NUM_COLS + col]",
            font_size=20, bold=True, color=RGBColor(0x1E, 0x40, 0xAF),
            font_name="Courier New")

bullets2 = [
    "\u2022  Elements in the same row are contiguous in memory (addresses differ by 1)",
    "\u2022  Elements in the same column are far apart (addresses differ by NUM_COLS)",
    "\u2022  Accessing a column = strided access = multiple DRAM transactions",
]
add_bullet_slide_content(slide, Inches(0.8), Inches(4.9), Inches(11.5), Inches(2.0),
                          bullets2, font_size=18, spacing=Pt(12))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6 – Visualizing the Access Pattern (Uncoalesced)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, *WHITE_BG)
add_section_header(slide, "Uncoalesced Access Pattern (Naive Kernel)")

add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
            "threadIdx.x \u2192 row (i)  |  A warp of 4 threads accesses 4 different rows",
            font_size=18, color=RGBColor(0x64, 0x74, 0x8B))

# Matrix A visualization
add_textbox(slide, Inches(0.8), Inches(2.1), Inches(3), Inches(0.4),
            "Matrix A  (reading column k=0)", font_size=14, bold=True,
            color=RGBColor(0xEF, 0x44, 0x44))

thread_colors = [
    RGBColor(0xEF, 0x44, 0x44),
    RGBColor(0xF9, 0x73, 0x16),
    RGBColor(0xEA, 0xB3, 0x08),
    RGBColor(0x84, 0xCC, 0x16),
]

# A matrix 6x6 with 4 highlighted cells in column 0
for r in range(6):
    for c in range(6):
        x = Inches(0.8 + c * 0.5)
        y = Inches(2.5 + r * 0.5)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x, y, Inches(0.45), Inches(0.45))
        shape.fill.solid()
        if c == 0 and r < 4:
            shape.fill.fore_color.rgb = thread_colors[r]
            tf = shape.text_frame
            tf.margin_left = Pt(0); tf.margin_right = Pt(0)
            p = tf.paragraphs[0]
            p.text = f"T{r}"
            p.font.size = Pt(9); p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER
        else:
            shape.fill.fore_color.rgb = RGBColor(0xFE, 0xE2, 0xE2)
        shape.line.fill.background()

# Arrow down showing strided access in DRAM
add_textbox(slide, Inches(4.2), Inches(3.2), Inches(1), Inches(0.5),
            "\u27a1", font_size=36, color=RGBColor(0x64, 0x74, 0x8B),
            alignment=PP_ALIGN.CENTER)

# DRAM ribbon for A
add_textbox(slide, Inches(5.5), Inches(2.1), Inches(7), Inches(0.4),
            "A in DRAM (row-major) — threads touch scattered locations",
            font_size=14, bold=True, color=RGBColor(0xEF, 0x44, 0x44))

# 36 cells for a 6x6 matrix, highlight positions 0, 6, 12, 18
for i in range(36):
    x = Inches(5.5 + i * 0.2)
    y = Inches(2.6)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x, y, Inches(0.18), Inches(0.5))
    shape.fill.solid()
    if i == 0:
        shape.fill.fore_color.rgb = thread_colors[0]
    elif i == 6:
        shape.fill.fore_color.rgb = thread_colors[1]
    elif i == 12:
        shape.fill.fore_color.rgb = thread_colors[2]
    elif i == 18:
        shape.fill.fore_color.rgb = thread_colors[3]
    else:
        shape.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.fill.background()

# Labels for highlighted DRAM positions
for idx, pos in enumerate([0, 6, 12, 18]):
    add_textbox(slide, Inches(5.5 + pos * 0.2 - 0.05), Inches(3.2),
                Inches(0.4), Inches(0.3),
                f"T{idx}", font_size=9, bold=True,
                color=thread_colors[idx], alignment=PP_ALIGN.CENTER)

# Strided pattern explanation
add_colored_rect(slide, Inches(5.5), Inches(3.7), Inches(7), Inches(0.7),
                  RGBColor(0xFE, 0xF2, 0xF2),
                  "Stride = K = 6 floats apart \u2192 each access hits a different "
                  "32-byte cache line \u2192 4 DRAM transactions for 4 floats!",
                  font_size=14, font_color=RGBColor(0x99, 0x1B, 0x1B))

# Bottom: efficiency comparison
add_textbox(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(0.4),
            "Bus Efficiency: 25%", font_size=28, bold=True,
            color=RGBColor(0xDC, 0x26, 0x26))

add_textbox(slide, Inches(0.8), Inches(5.6), Inches(10), Inches(0.8),
            "Only 4 useful bytes out of every 32-byte transaction are used.\n"
            "We pay for 128 bytes of DRAM bandwidth but only need 16 bytes.",
            font_size=16, color=RGBColor(0x64, 0x74, 0x8B))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7 – The Fix: Coalesced Access
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, *WHITE_BG)
add_section_header(slide, "The Fix: Swap the Mapping")

add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
            "One line change — massive performance impact",
            font_size=20, bold=True, color=RGBColor(0x16, 0x65, 0x34))

# Left: before (bad)
add_colored_rect(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.0),
                  RGBColor(0xFE, 0xF2, 0xF2))
add_textbox(slide, Inches(1.0), Inches(2.3), Inches(5), Inches(0.4),
            "\u274c  Uncoalesced (Naive)", font_size=18, bold=True,
            color=RGBColor(0xDC, 0x26, 0x26))
add_code_box(slide, Inches(1.0), Inches(2.8), Inches(5), Inches(1.1),
             "// threadIdx.x -> row\n"
             "i = blockIdx.x * blockDim.x + threadIdx.x;\n"
             "j = blockIdx.y * blockDim.y + threadIdx.y;",
             font_size=14)

# Right: after (good)
add_colored_rect(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(2.0),
                  RGBColor(0xF0, 0xFD, 0xF4))
add_textbox(slide, Inches(7.2), Inches(2.3), Inches(5), Inches(0.4),
            "\u2705  Coalesced (Fixed)", font_size=18, bold=True,
            color=RGBColor(0x16, 0x65, 0x34))
add_code_box(slide, Inches(7.2), Inches(2.8), Inches(5), Inches(1.1),
             "// threadIdx.x -> column\n"
             "j = blockIdx.x * blockDim.x + threadIdx.x;\n"
             "i = blockIdx.y * blockDim.y + threadIdx.y;",
             font_size=14)

# Access pattern comparison
add_textbox(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(0.3),
            "Matrix B access: B[k*N + j]", font_size=16, bold=True,
            color=RGBColor(0x1E, 0x3A, 0x5F))

# DRAM ribbon - Uncoalesced B access
add_textbox(slide, Inches(0.8), Inches(5.1), Inches(3), Inches(0.25),
            "Naive: j is constant \u2192 broadcast (1 value)", font_size=12,
            color=RGBColor(0xDC, 0x26, 0x26))

for i in range(32):
    x = Inches(0.8 + i * 0.17)
    y = Inches(5.4)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x, y, Inches(0.15), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF9, 0x73, 0x16) if i == 3 else RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.fill.background()

# DRAM ribbon - Coalesced B access
add_textbox(slide, Inches(7), Inches(5.1), Inches(5), Inches(0.25),
            "Fixed: j varies with threadIdx.x \u2192 consecutive addresses",
            font_size=12, color=RGBColor(0x16, 0x65, 0x34))

for i in range(32):
    x = Inches(7 + i * 0.17)
    y = Inches(5.4)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x, y, Inches(0.15), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x22, 0xC5, 0x5E) if i < 4 else RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.fill.background()

# Bottom result
add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.6),
            "Result: Bus efficiency goes from 25% \u2192 100% for B reads. "
            "One-line fix, huge speedup.",
            font_size=20, bold=True, color=RGBColor(0x16, 0x65, 0x34))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8 – DRAM Transaction Deep Dive
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, *WHITE_BG)
add_section_header(slide, "DRAM Transactions: The Numbers")

# Transaction comparison table
headers = ["", "Uncoalesced (Naive)", "Coalesced (Fixed)"]
rows_data = [
    ["A read pattern", "Strided (each thread hits\na different cache line)", "Broadcast (all threads read\nsame element \u2192 1 transaction)"],
    ["B read pattern", "Broadcast (all threads read\nsame element \u2192 1 transaction)", "Coalesced (adjacent threads\nread adjacent addresses \u2192 1 txn)"],
    ["C write pattern", "Strided", "Coalesced"],
    ["Transactions per\nwarp per k-step", "~5 (strided A + broadcast B)", "2 (broadcast A + coalesced B)"],
    ["Useful bytes /\ntotal bytes moved", "16B / 128B = 12.5%", "16B / 64B = 25-50%"],
    ["Bus Efficiency", "25%", "100%"],
]

# Table header
for ci, h in enumerate(headers):
    x = Inches(0.8 + ci * 4)
    w = Inches(3.8) if ci > 0 else Inches(2.5)
    box = add_colored_rect(slide, x, Inches(1.5), w, Inches(0.5),
                            RGBColor(0x1E, 0x29, 0x3B), h, font_size=13,
                            font_color=RGBColor(0xFF, 0xFF, 0xFF))

for ri, row in enumerate(rows_data):
    y = Inches(2.05 + ri * 0.78)
    bg = RGBColor(0xF8, 0xFA, 0xFC) if ri % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
    for ci, cell_text in enumerate(row):
        x = Inches(0.8 + ci * 4)
        w = Inches(3.8) if ci > 0 else Inches(2.5)
        if ri == len(rows_data) - 1:
            # Last row - highlight
            if ci == 1:
                bg_c = RGBColor(0xFE, 0xE2, 0xE2)
                fc = RGBColor(0x99, 0x1B, 0x1B)
            elif ci == 2:
                bg_c = RGBColor(0xDC, 0xFC, 0xE7)
                fc = RGBColor(0x16, 0x65, 0x34)
            else:
                bg_c = bg
                fc = RGBColor(0x33, 0x33, 0x33)
            add_colored_rect(slide, x, y, w, Inches(0.7), bg_c, cell_text,
                              font_size=12, font_color=fc)
        else:
            box = add_colored_rect(slide, x, y, w, Inches(0.7), bg)
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.15)
            tf.margin_top = Inches(0.05)
            p = tf.paragraphs[0]
            p.text = cell_text
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.font.name = "Calibri"

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 9 – Why This Matters / Performance Impact
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, *WHITE_BG)
add_section_header(slide, "Why Does This Matter?")

# Left - The problem
add_colored_rect(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.5),
                  RGBColor(0xFE, 0xF2, 0xF2))

add_textbox(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.4),
            "The Naive Kernel Bottleneck", font_size=20, bold=True,
            color=RGBColor(0xDC, 0x26, 0x26))

left_items = [
    "\u2022  Strided access \u2192 each thread triggers its own\n   DRAM transaction (32B each)",
    "\u2022  A 32-thread warp needing 32 floats (128B)\n   may cause up to 32 transactions = 1024B moved",
    "\u2022  That is 8\u00d7 more data moved than needed",
    "\u2022  Memory bus becomes the bottleneck\n   long before compute is saturated",
    "\u2022  GPU cores sit idle waiting for data",
]
add_bullet_slide_content(slide, Inches(1.0), Inches(2.3), Inches(5), Inches(3.5),
                          left_items, font_size=16, spacing=Pt(14),
                          color=RGBColor(0x7F, 0x1D, 0x1D))

# Right - The fix benefit
add_colored_rect(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(4.5),
                  RGBColor(0xF0, 0xFD, 0xF4))

add_textbox(slide, Inches(7.2), Inches(1.7), Inches(5), Inches(0.4),
            "After Coalescing", font_size=20, bold=True,
            color=RGBColor(0x16, 0x65, 0x34))

right_items = [
    "\u2022  Adjacent threads \u2192 adjacent addresses\n   \u2192 single DRAM transaction",
    "\u2022  32 threads needing 32 floats = 128B\n   = only 4 transactions (perfectly packed)",
    "\u2022  4\u00d7 to 8\u00d7 bandwidth improvement",
    "\u2022  This alone can mean 2\u00d7-10\u00d7 speedup\n   depending on the kernel",
    "\u2022  Still just the first step —\n   much more optimization ahead!",
]
add_bullet_slide_content(slide, Inches(7.2), Inches(2.3), Inches(5), Inches(3.5),
                          right_items, font_size=16, spacing=Pt(14),
                          color=RGBColor(0x14, 0x53, 0x2D))

# Bottom banner
add_colored_rect(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8),
                  RGBColor(0x1E, 0x40, 0xAF),
                  "Rule #1 of GPU Programming: Always coalesce your memory accesses",
                  font_size=20, font_color=RGBColor(0xFF, 0xFF, 0xFF))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 10 – Road Ahead
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, *DARK_BG)

add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
            "The Optimization Roadmap", font_size=36, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF))

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.8), Inches(1.5), Inches(2), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0x3B, 0x82, 0xF6)
line.line.fill.background()

stages = [
    ("0", "Naive (Uncoalesced)", "The anti-pattern", RGBColor(0xEF, 0x44, 0x44), True),
    ("1", "Coalesced Access", "Swap i/j mapping", RGBColor(0x22, 0xC5, 0x5E), True),
    ("2", "Shared Memory Tiling", "Reduce global memory traffic", RGBColor(0x3B, 0x82, 0xF6), False),
    ("3", "Register Tiling (1D)", "More work per thread", RGBColor(0x8B, 0x5C, 0xF6), False),
    ("4", "Register Tiling (2D)", "8\u00d78 output per thread", RGBColor(0xEC, 0x48, 0x99), False),
    ("5", "Vectorized (float4)", "128-bit wide loads", RGBColor(0xF5, 0x9E, 0x0B), False),
    ("6", "Tensor Cores (WMMA)", "Hardware matrix units", RGBColor(0x06, 0xB6, 0xD4), False),
]

for idx, (num, title, desc, color, done) in enumerate(stages):
    x = Inches(0.8 + (idx % 4) * 3.1)
    y = Inches(2.2 + (idx // 4) * 2.3)

    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, y, Inches(2.8), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    # Stage number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x + Inches(0.15), y + Inches(0.15),
                                     Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Check mark for completed
    if done:
        add_textbox(slide, x + Inches(2.2), y + Inches(0.1), Inches(0.5), Inches(0.4),
                    "\u2713", font_size=20, bold=True, color=RGBColor(0x22, 0xC5, 0x5E))

    # Title
    add_textbox(slide, x + Inches(0.15), y + Inches(0.7), Inches(2.5), Inches(0.5),
                title, font_size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    # Description
    add_textbox(slide, x + Inches(0.15), y + Inches(1.2), Inches(2.5), Inches(0.5),
                desc, font_size=13, color=RGBColor(0x94, 0xA3, 0xB8))

add_textbox(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
            "We just covered stages 0 & 1  \u2014  stay tuned for the rest of the journey!",
            font_size=16, color=RGBColor(0x94, 0xA3, 0xB8))

# ── Save ─────────────────────────────────────────────────────────────────────
output_path = "matmul/matmul_presentation.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
